"""
Presentation Themes — Config-Driven Visual Theme Definitions

Each theme is a dataclass containing color palettes, font pairings, and a
background renderer. Themes are fully self-contained and require no external
image files; all visuals are drawn programmatically via python-pptx shapes.
"""

from dataclasses import dataclass, field
from typing import Callable, Optional
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


@dataclass
class PresentationTheme:
    """A complete visual theme for presentation generation."""
    name: str
    display_name: str

    # Color palette
    primary_color: RGBColor
    secondary_color: RGBColor
    accent_color: RGBColor
    bg_color: RGBColor
    title_font_color: RGBColor
    body_font_color: RGBColor
    subtitle_font_color: RGBColor

    # Font pairing
    title_font: str = "Calibri"
    body_font: str = "Calibri Light"

    # Font sizes
    cover_title_size: Pt = field(default_factory=lambda: Pt(36))
    cover_subtitle_size: Pt = field(default_factory=lambda: Pt(16))
    slide_title_size: Pt = field(default_factory=lambda: Pt(28))
    body_font_size: Pt = field(default_factory=lambda: Pt(18))

    # Bullet character
    bullet_char: str = "•"


# ---------------------------------------------------------------------------
# Background renderers — each draws abstract shapes on a slide
# ---------------------------------------------------------------------------

def _render_gradient_bar(slide, prs, theme: PresentationTheme, is_title_slide: bool = False):
    """Draws a full-width colored header bar and a thin accent stripe."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if is_title_slide:
        # Full-height dark background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.primary_color
        bg.line.fill.background()
        # Accent bar at bottom quarter
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, int(slide_h * 0.78),
            slide_w, int(slide_h * 0.04)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = theme.accent_color
        accent.line.fill.background()
    else:
        # Light background fill
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.bg_color
        bg.line.fill.background()
        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = theme.primary_color
        header.line.fill.background()
        # Thin accent stripe below header
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.2), slide_w, Inches(0.06)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = theme.accent_color
        stripe.line.fill.background()


def _render_sidebar_accent(slide, prs, theme: PresentationTheme, is_title_slide: bool = False):
    """Draws a left sidebar accent band and optional corner shape."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if is_title_slide:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.bg_color
        bg.line.fill.background()
        # Wide left sidebar
        sidebar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(1.5), slide_h
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = theme.primary_color
        sidebar.line.fill.background()
        # Small accent square in bottom-left
        sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.2), int(slide_h - Inches(1.0)),
            Inches(1.1), Inches(0.8)
        )
        sq.fill.solid()
        sq.fill.fore_color.rgb = theme.accent_color
        sq.line.fill.background()
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.bg_color
        bg.line.fill.background()
        # Narrow left accent bar
        sidebar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), slide_h
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = theme.primary_color
        sidebar.line.fill.background()
        # Header area (light secondary)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), 0,
            int(slide_w - Inches(0.4)), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = theme.secondary_color
        header.line.fill.background()


def _render_top_bottom_bands(slide, prs, theme: PresentationTheme, is_title_slide: bool = False):
    """Draws colored bands at top and bottom of the slide."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Full background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.bg_color
    bg.line.fill.background()

    if is_title_slide:
        # Top band (60% height)
        top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_w, int(slide_h * 0.6)
        )
        top.fill.solid()
        top.fill.fore_color.rgb = theme.primary_color
        top.line.fill.background()
        # Accent stripe
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, int(slide_h * 0.6),
            slide_w, Inches(0.08)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = theme.accent_color
        stripe.line.fill.background()
    else:
        # Top band
        top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(1.3)
        )
        top.fill.solid()
        top.fill.fore_color.rgb = theme.primary_color
        top.line.fill.background()
        # Bottom stripe
        bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, int(slide_h - Inches(0.35)),
            slide_w, Inches(0.35)
        )
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = theme.secondary_color
        bottom.line.fill.background()


def _render_minimal_header(slide, prs, theme: PresentationTheme, is_title_slide: bool = False):
    """Clean minimal style — thin accent line only, mostly white space."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.bg_color
    bg.line.fill.background()

    if is_title_slide:
        # Centered accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(slide_w * 0.3), int(slide_h * 0.55),
            int(slide_w * 0.4), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = theme.accent_color
        line.line.fill.background()
    else:
        # Top thin line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.15),
            slide_w, Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = theme.accent_color
        line.line.fill.background()


def _render_corner_blocks(slide, prs, theme: PresentationTheme, is_title_slide: bool = False):
    """Geometric corner accent blocks for a modern feel."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.bg_color
    bg.line.fill.background()

    if is_title_slide:
        # Large top-right block
        tr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(slide_w - Inches(3.5)), 0,
            Inches(3.5), Inches(2.5)
        )
        tr.fill.solid()
        tr.fill.fore_color.rgb = theme.primary_color
        tr.line.fill.background()
        # Small bottom-left block
        bl = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, int(slide_h - Inches(2.0)),
            Inches(2.5), Inches(2.0)
        )
        bl.fill.solid()
        bl.fill.fore_color.rgb = theme.accent_color
        bl.line.fill.background()
    else:
        # Top header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = theme.primary_color
        header.line.fill.background()
        # Small accent square bottom-right
        sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(slide_w - Inches(0.6)), int(slide_h - Inches(0.6)),
            Inches(0.6), Inches(0.6)
        )
        sq.fill.solid()
        sq.fill.fore_color.rgb = theme.accent_color
        sq.line.fill.background()


# ---------------------------------------------------------------------------
# Theme Definitions
# ---------------------------------------------------------------------------

THEMES = {
    "corporate_navy": PresentationTheme(
        name="corporate_navy",
        display_name="Corporate Navy",
        primary_color=RGBColor(0x1A, 0x1C, 0x2C),
        secondary_color=RGBColor(0x42, 0x47, 0x69),
        accent_color=RGBColor(0x70, 0x77, 0xA1),
        bg_color=RGBColor(0xF8, 0xFA, 0xFC),
        title_font_color=RGBColor(0xFF, 0xFF, 0xFF),
        body_font_color=RGBColor(0x1E, 0x1E, 0x1E),
        subtitle_font_color=RGBColor(0xC0, 0xC4, 0xD4),
        title_font="Calibri",
        body_font="Calibri Light",
    ),
    "emerald_executive": PresentationTheme(
        name="emerald_executive",
        display_name="Emerald Executive",
        primary_color=RGBColor(0x0B, 0x3D, 0x2E),
        secondary_color=RGBColor(0x14, 0x5A, 0x46),
        accent_color=RGBColor(0xD4, 0xA5, 0x37),
        bg_color=RGBColor(0xFA, 0xFA, 0xF5),
        title_font_color=RGBColor(0xFF, 0xFF, 0xFF),
        body_font_color=RGBColor(0x1C, 0x1C, 0x1C),
        subtitle_font_color=RGBColor(0xD4, 0xA5, 0x37),
        title_font="Georgia",
        body_font="Calibri",
    ),
    "slate_minimal": PresentationTheme(
        name="slate_minimal",
        display_name="Slate Minimal",
        primary_color=RGBColor(0x33, 0x41, 0x55),
        secondary_color=RGBColor(0x64, 0x74, 0x8B),
        accent_color=RGBColor(0x38, 0xBD, 0xF8),
        bg_color=RGBColor(0xFF, 0xFF, 0xFF),
        title_font_color=RGBColor(0x33, 0x41, 0x55),
        body_font_color=RGBColor(0x33, 0x41, 0x55),
        subtitle_font_color=RGBColor(0x64, 0x74, 0x8B),
        title_font="Calibri",
        body_font="Calibri Light",
    ),
    "crimson_bold": PresentationTheme(
        name="crimson_bold",
        display_name="Crimson Bold",
        primary_color=RGBColor(0x7F, 0x1D, 0x1D),
        secondary_color=RGBColor(0x4A, 0x14, 0x14),
        accent_color=RGBColor(0xF5, 0x9E, 0x0B),
        bg_color=RGBColor(0xFE, 0xF9, 0xF2),
        title_font_color=RGBColor(0xFF, 0xFF, 0xFF),
        body_font_color=RGBColor(0x1E, 0x1E, 0x1E),
        subtitle_font_color=RGBColor(0xF5, 0x9E, 0x0B),
        title_font="Calibri",
        body_font="Calibri",
    ),
    "azure_modern": PresentationTheme(
        name="azure_modern",
        display_name="Azure Modern",
        primary_color=RGBColor(0x1E, 0x40, 0xAF),
        secondary_color=RGBColor(0x31, 0x30, 0x64),
        accent_color=RGBColor(0x06, 0xB6, 0xD4),
        bg_color=RGBColor(0xF0, 0xF9, 0xFF),
        title_font_color=RGBColor(0xFF, 0xFF, 0xFF),
        body_font_color=RGBColor(0x1E, 0x29, 0x3B),
        subtitle_font_color=RGBColor(0x93, 0xC5, 0xFD),
        title_font="Calibri",
        body_font="Calibri Light",
    ),
}

# Map each theme to its background renderer
_RENDERERS = {
    "corporate_navy": _render_gradient_bar,
    "emerald_executive": _render_sidebar_accent,
    "slate_minimal": _render_minimal_header,
    "crimson_bold": _render_top_bottom_bands,
    "azure_modern": _render_corner_blocks,
}

DEFAULT_THEME = "corporate_navy"


def get_theme(name: str) -> PresentationTheme:
    """Retrieve a theme by internal name, falling back to the default."""
    return THEMES.get(name, THEMES[DEFAULT_THEME])


def get_theme_names() -> list:
    """Returns a list of human-readable theme display names."""
    return [t.display_name for t in THEMES.values()]


def get_theme_by_display_name(display_name: str) -> PresentationTheme:
    """Retrieve a theme by its user-facing display name."""
    for t in THEMES.values():
        if t.display_name == display_name:
            return t
    return THEMES[DEFAULT_THEME]


def render_background(slide, prs, theme: PresentationTheme, is_title_slide: bool = False):
    """Renders the theme's background shapes onto a slide."""
    renderer = _RENDERERS.get(theme.name, _render_gradient_bar)
    renderer(slide, prs, theme, is_title_slide)
