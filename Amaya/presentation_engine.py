import logging
import re
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import config
from presentation_themes import (
    PresentationTheme, get_theme, get_theme_by_display_name,
    render_background, DEFAULT_THEME
)

logger = logging.getLogger("PresentationEngine")


class PresentationEngine:
    """
    Professional PowerPoint Generation Engine.

    Produces themed, visually styled slides with proper title/content layouts,
    background shapes, and styled bullet points. Themes are config-driven via
    presentation_themes.py.
    """

    BULLETS_PER_SLIDE = 6  # Max bullets before auto-splitting into continuation

    def __init__(self):
        self.output_root = Path(config.OUTPUT_ROOT)

    # ------------------------------------------------------------------
    # Slide builders
    # ------------------------------------------------------------------

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str, theme: PresentationTheme):
        """Creates a themed cover/title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        render_background(slide, prs, theme, is_title_slide=True)

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Title text box
        left = Inches(1.0)
        top = int(slide_h * 0.32)
        width = int(slide_w - Inches(2.0))
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = theme.cover_title_size
        p.font.bold = True
        p.font.name = theme.title_font
        p.font.color.rgb = theme.title_font_color
        p.alignment = PP_ALIGN.LEFT

        # Subtitle text box
        sub_top = int(slide_h * 0.52)
        sub_height = Inches(0.8)
        subBox = slide.shapes.add_textbox(left, sub_top, width, sub_height)
        sf = subBox.text_frame
        sf.word_wrap = True
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = theme.cover_subtitle_size
        sp.font.name = theme.body_font
        sp.font.color.rgb = theme.subtitle_font_color
        sp.alignment = PP_ALIGN.LEFT

    def _add_content_slide(self, prs: Presentation, title: str,
                           bullets: List[str], theme: PresentationTheme):
        """Creates a single themed content slide with title and bullet points."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        render_background(slide, prs, theme, is_title_slide=False)

        slide_w = prs.slide_width

        # Title text box (overlays the header bar area)
        title_left = Inches(0.8)
        title_top = Inches(0.2)
        title_width = int(slide_w - Inches(1.6))
        title_height = Inches(0.9)
        title_box = slide.shapes.add_textbox(title_left, title_top,
                                             title_width, title_height)
        ttf = title_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = title
        tp.font.size = theme.slide_title_size
        tp.font.bold = True
        tp.font.name = theme.title_font
        tp.font.color.rgb = theme.title_font_color
        tp.alignment = PP_ALIGN.LEFT

        # Content text box (below header)
        body_left = Inches(0.8)
        body_top = Inches(1.5)
        body_width = int(slide_w - Inches(1.6))
        body_height = Inches(5.0)
        body_box = slide.shapes.add_textbox(body_left, body_top,
                                            body_width, body_height)
        btf = body_box.text_frame
        btf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            # Clean any residual markdown
            clean = re.sub(r'^\*\*|\*\*$', '', bullet).strip()
            p.text = f"{theme.bullet_char}  {clean}"
            p.font.size = theme.body_font_size
            p.font.name = theme.body_font
            p.font.color.rgb = theme.body_font_color
            p.space_after = Pt(8)
            p.level = 0

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def create_presentation(self, title: str, subtitle: str,
                            slides_data: List[Dict[str, Any]],
                            output_path: Path,
                            theme_name: str = None) -> str:
        """
        Constructs a professionally themed .pptx file.

        Args:
            title: Cover slide title.
            subtitle: Cover slide subtitle.
            slides_data: List of {"title": str, "content": [str, ...]}.
            output_path: Where to save the .pptx file.
            theme_name: Internal name or display name of the theme.
                        Falls back to DEFAULT_THEME if invalid/None.

        Returns:
            str: Path to the generated .pptx file.
        """
        # Resolve theme
        if theme_name:
            theme = get_theme(theme_name)
            # If internal name lookup returned default but user passed a display name, try that
            if theme.name == DEFAULT_THEME and theme_name != DEFAULT_THEME:
                theme = get_theme_by_display_name(theme_name)
        else:
            theme = get_theme(DEFAULT_THEME)

        prs = Presentation()

        # 1. Title slide
        self._add_title_slide(prs, title, subtitle, theme)

        # 2. Content slides (with auto-continuation for overflow)
        for data in slides_data:
            slide_title = data.get("title", "Strategic Insight")
            content = data.get("content", [])

            if not content:
                # Empty slide — still render it with the title
                self._add_content_slide(prs, slide_title, [], theme)
                continue

            # Split into chunks of BULLETS_PER_SLIDE
            for chunk_idx in range(0, len(content), self.BULLETS_PER_SLIDE):
                chunk = content[chunk_idx:chunk_idx + self.BULLETS_PER_SLIDE]
                if chunk_idx == 0:
                    self._add_content_slide(prs, slide_title, chunk, theme)
                else:
                    cont_title = f"{slide_title} (cont'd)"
                    self._add_content_slide(prs, cont_title, chunk, theme)

        # 3. Save
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info(f"Themed presentation generated: {output_path}")
        return str(output_path)

    def parse_llm_summary(self, llm_text: str) -> List[Dict[str, Any]]:
        """
        Robustly parses LLM output into structured slide data.
        Discards conversational filler and cleans up unwanted labels.
        """
        import re
        slides = []
        current_slide = None

        # 1. Split into lines and find the first SLIDE: tag
        lines = llm_text.split("\n")
        start_index = -1
        for i, line in enumerate(lines):
            if re.search(r'^\s*(?:\*\*|###)?\s*SLIDE\b', line, re.IGNORECASE):
                start_index = i
                break

        if start_index == -1:
            logger.warning("No SLIDE: tags found in LLM output.")
            return []

        # 2. Process only from the first SLIDE: tag onwards
        for line in lines[start_index:]:
            line = line.strip()
            if not line: continue

            # Match SLIDE: Title
            header_match = re.search(r'(?:\*\*|###)?\s*SLIDE\s*(?:[0-9]+)?\s*[:\-]*\s*(?:\*\*)?(.*)', line, re.IGNORECASE)
            if header_match:
                if current_slide:
                    slides.append(current_slide)

                title = header_match.group(1).strip()
                # Clean up any trailing markdown or image hints like "(Image: ...)"
                title = re.sub(r'\(Image:.*?\)', '', title, flags=re.IGNORECASE)
                title = re.sub(r'^\*\*|\*\*$', '', title).strip()
                if not title: title = "Insight"

                current_slide = {"title": title, "content": []}
                continue

            if not current_slide: continue

            # Match Bullet Points or any informative lines
            # Clean up labels like "Headline:", "Impact:", "Key Benefit:", etc.
            clean_line = re.sub(r'^(?:[-*]\s*)?(?:Headline|Impact|Key Benefit|Subtitle|Briefly state|Result|Goal|Bullet Points)\s*[:\-]*\s*', '', line, flags=re.IGNORECASE).strip()
            # Strip leading bullet hyphens or asterisks
            clean_line = re.sub(r'^[-*]\s*', '', clean_line).strip()
            # Remove any remaining Markdown bolding
            clean_line = re.sub(r'^\*\*|\*\*$', '', clean_line).strip()
            # Remove image hints
            clean_line = re.sub(r'\(Image:.*?\)', '', clean_line, flags=re.IGNORECASE).strip()
            # Remove separator lines like "--"
            if clean_line and clean_line != "--":
                current_slide["content"].append(clean_line)

        if current_slide:
            slides.append(current_slide)

        return slides
